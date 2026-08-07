"""PowerPoint parsing with explicit ``python_pptx`` and ``liteparse`` providers."""

from __future__ import annotations

import asyncio
import logging
from pathlib import Path
import shutil
import subprocess
import tempfile
from typing import Any, Optional, Protocol, TYPE_CHECKING

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    Presentation = None
    MSO_SHAPE_TYPE = None

from .asset_reference import AssetReferenceMode, prepare_markdown_asset_references
from .base_document_service import ASSET_STAGE_NAMES, BaseDocumentService
from .document_artifact_models import DocumentAsset, MarkdownArtifact
from .document_artifact_writer import DocumentArtifactWriter
from .document_asset_publisher import (
    AftsDocumentAssetPublisher,
    DocumentAssetPublisher,
    NoOpDocumentAssetPublisher,
)
from .document_parse_logging import DocumentParseLogger
from .markdown_assembler import MarkdownAssembler, PlaceholderMarkdownAssembler
from .paths import DOCUMENT_PARSE_WORKSPACE

if TYPE_CHECKING:
    from services.afts_service import AftsService


logger = logging.getLogger(__name__)

PPTX_PARSE_PROVIDERS = ("python_pptx", "liteparse")
DEFAULT_PPTX_PARSE_PROVIDER = "python_pptx"
_IMAGE_CONTENT_TYPES = {
    "image/png": ".png",
    "image/jpeg": ".jpg",
    "image/gif": ".gif",
    "image/bmp": ".bmp",
    "image/tiff": ".tiff",
    "image/webp": ".webp",
}


class PptTextProvider(Protocol):
    """Produces ordered Markdown text for a presentation."""

    name: str

    async def extract_markdown(self, file_path: Path) -> str:
        """Extract presentation content as Markdown."""


class PythonPptxProvider:
    """Extract native slide text and tables with python-pptx."""

    name = "python_pptx"

    async def extract_markdown(self, file_path: Path) -> str:
        return await asyncio.to_thread(self._extract_markdown, file_path)

    def _extract_markdown(self, file_path: Path) -> str:
        presentation = _load_presentation(file_path)
        parts: list[str] = []
        for slide_number, slide in enumerate(presentation.slides, start=1):
            parts.append(f"## 幻灯片 {slide_number}\n\n")
            slide_has_content = False
            for shape in _iter_shapes(slide.shapes):
                if getattr(shape, "has_table", False):
                    table_markdown = self._table_to_markdown(shape.table)
                    if table_markdown:
                        parts.append(table_markdown)
                        slide_has_content = True
                    continue
                text = str(getattr(shape, "text", "") or "").strip()
                if text:
                    parts.append(f"{text}\n\n")
                    slide_has_content = True
            if not slide_has_content:
                parts.append("_本页无可提取文本_\n\n")
        return "".join(parts).rstrip()

    @staticmethod
    def _table_to_markdown(table: Any) -> str:
        rows = [
            [_escape_table_cell(getattr(cell, "text", "")) for cell in row.cells]
            for row in table.rows
        ]
        if not rows:
            return ""
        column_count = max(len(row) for row in rows)
        normalized_rows = [row + [""] * (column_count - len(row)) for row in rows]
        header = normalized_rows[0]
        lines = [
            "| " + " | ".join(header) + " |",
            "| " + " | ".join("---" for _ in range(column_count)) + " |",
        ]
        lines.extend("| " + " | ".join(row) + " |" for row in normalized_rows[1:])
        return "\n".join(lines) + "\n\n"


class LiteParseProvider:
    """Extract PPT/PPTX text through the installed LiteParse CLI."""

    name = "liteparse"

    def __init__(self, *, cli_path: str = "", ocr_enabled: bool = False) -> None:
        self._cli_path = str(cli_path or "").strip()
        self._ocr_enabled = ocr_enabled

    async def extract_markdown(self, file_path: Path) -> str:
        cli_command = self._resolve_cli_command()
        with tempfile.TemporaryDirectory(prefix="filex-liteparse-pptx-") as temp_dir:
            output_path = Path(temp_dir) / "output.txt"
            command = [
                *cli_command,
                "parse",
                str(file_path),
                "--format",
                "text",
                "-q",
                "-o",
                str(output_path),
            ]
            if not self._ocr_enabled:
                command.append("--no-ocr")
            result = await asyncio.to_thread(
                subprocess.run,
                command,
                capture_output=True,
                text=True,
                check=False,
            )
            if result.returncode != 0:
                error = (result.stderr or result.stdout or "").strip()
                raise RuntimeError(
                    "LiteParse PPTX parse failed: "
                    + (error or f"exit code {result.returncode}")
                )
            markdown_text = (
                output_path.read_text(encoding="utf-8").strip()
                if output_path.exists()
                else ""
            )
        if not markdown_text:
            raise RuntimeError("LiteParse PPTX parse result is empty")
        return markdown_text

    def _resolve_cli_command(self) -> list[str]:
        if self._cli_path:
            configured = Path(self._cli_path).expanduser()
            if configured.is_file():
                return [str(configured)]
            raise RuntimeError(f"Configured LiteParse CLI does not exist: {configured}")
        for candidate in ("liteparse", "lit"):
            executable = shutil.which(candidate)
            if executable:
                return [executable]
        raise RuntimeError(
            "LiteParse provider requires the `liteparse` or `lit` CLI in PATH"
        )


class PptDocumentService(BaseDocumentService):
    """Parse PPT/PPTX with an explicit text provider and shared image extraction."""

    _stage_names = ASSET_STAGE_NAMES
    _default_suffix = "pptx"
    _empty_error_message = "PPTX 解析结果为空"

    def __init__(
        self,
        *,
        env_content: Optional[dict[str, Any]] = None,
        asset_reference_mode: AssetReferenceMode = "remote_id",
        text_provider: PptTextProvider | None = None,
        markdown_assembler: MarkdownAssembler | None = None,
        artifact_writer: DocumentArtifactWriter | None = None,
    ) -> None:
        super().__init__(artifact_writer=artifact_writer)
        self._env_content = env_content or {}
        self._asset_reference_mode = asset_reference_mode
        self._text_provider = text_provider
        self._markdown_assembler = markdown_assembler or PlaceholderMarkdownAssembler()

    async def _build_artifact(
        self,
        *,
        file_path: Path,
        task_id: str,
        source_file_name: str,
        afts_service: Optional["AftsService"],
        stage_logger: DocumentParseLogger,
    ) -> MarkdownArtifact:
        provider = self._resolve_text_provider()
        output_dir = DOCUMENT_PARSE_WORKSPACE / task_id
        output_dir.mkdir(parents=True, exist_ok=True)

        with stage_logger.stage("content_extract", provider=provider.name):
            markdown_text = await provider.extract_markdown(file_path)
            presentation_metrics = await asyncio.to_thread(
                self._inspect_presentation,
                file_path,
            )

        with stage_logger.stage("asset_extract", provider="python_pptx"):
            assets = await asyncio.to_thread(
                self._extract_image_assets,
                file_path,
                output_dir / f"{source_file_name}_images",
            )
            published_assets = await self._build_asset_publisher(
                afts_service
            ).publish_assets(assets)
            prepare_markdown_asset_references(
                published_assets,
                output_dir=output_dir,
                asset_reference_mode=self._asset_reference_mode,
            )
            self._validate_published_assets(published_assets)

        artifact = MarkdownArtifact(
            markdown_text=markdown_text,
            assets=published_assets,
            diagnostics={
                "task_id": task_id,
                "source_file_name": source_file_name,
                "provider": provider.name,
                "asset_count": len(published_assets),
                **presentation_metrics,
            },
        )
        with stage_logger.stage(
            "markdown_assemble",
            provider=provider.name,
            asset_count=len(published_assets),
        ):
            artifact.markdown_text = self._markdown_assembler.assemble(artifact)
        return artifact

    @staticmethod
    def _inspect_presentation(file_path: Path) -> dict[str, int]:
        presentation = _load_presentation(file_path)
        metrics = {
            "slide_count": len(presentation.slides),
            "text_box_count": 0,
            "table_count": 0,
            "speaker_note_count": 0,
            "empty_slide_count": 0,
        }
        for slide in presentation.slides:
            has_content = False
            for shape in _iter_shapes(slide.shapes):
                if getattr(shape, "has_table", False):
                    metrics["table_count"] += 1
                    has_content = True
                if str(getattr(shape, "text", "") or "").strip():
                    metrics["text_box_count"] += 1
                    has_content = True
                if _is_picture_shape(shape):
                    has_content = True
            notes_slide = getattr(slide, "notes_slide", None)
            if notes_slide is not None and any(
                str(getattr(shape, "text", "") or "").strip()
                for shape in notes_slide.shapes
            ):
                metrics["speaker_note_count"] += 1
            if not has_content:
                metrics["empty_slide_count"] += 1
        return metrics

    def _resolve_text_provider(self) -> PptTextProvider:
        if self._text_provider is not None:
            return self._text_provider
        provider_name = str(
            self._env_content.get("pptx_parse_provider")
            or self._env_content.get("ppt_parse_provider")
            or DEFAULT_PPTX_PARSE_PROVIDER
        ).strip().lower()
        if provider_name == "python_pptx":
            self._text_provider = PythonPptxProvider()
        elif provider_name == "liteparse":
            self._text_provider = LiteParseProvider(
                cli_path=str(self._env_content.get("liteparse_cli_path") or ""),
                ocr_enabled=_parse_bool(
                    self._env_content.get("pptx_liteparse_ocr_enabled"),
                    default=False,
                ),
            )
        else:
            raise ValueError(
                f"Unsupported PPTX parse provider: {provider_name}. "
                f"Expected one of: {', '.join(PPTX_PARSE_PROVIDERS)}"
            )
        return self._text_provider

    @staticmethod
    def _extract_image_assets(file_path: Path, output_dir: Path) -> list[DocumentAsset]:
        presentation = _load_presentation(file_path)
        assets: list[DocumentAsset] = []
        for slide_number, slide in enumerate(presentation.slides, start=1):
            for shape in _iter_shapes(slide.shapes):
                if not _is_picture_shape(shape):
                    continue
                image = _get_shape_image_or_none(
                    shape,
                    file_path=file_path,
                    slide_number=slide_number,
                )
                if image is None:
                    continue
                extension = _image_extension(image)
                if extension is None:
                    logger.warning(
                        "ppt_document_service skip unsupported image format | "
                        "file_path=%s slide_number=%s shape_name=%s",
                        file_path,
                        slide_number,
                        str(getattr(shape, "name", "") or ""),
                    )
                    continue
                mime_type = _image_content_type(image, extension)
                order = len(assets) + 1
                output_dir.mkdir(parents=True, exist_ok=True)
                local_path = output_dir / f"{file_path.stem}_slide_{slide_number}_img_{order}{extension}"
                local_path.write_bytes(image.blob)
                assets.append(
                    DocumentAsset(
                        asset_id=f"pptx_image_{order}",
                        kind="embedded_image",
                        local_path=local_path,
                        page_number=slide_number,
                        order=order,
                        meta={
                            "index": str(order),
                            "slide_number": slide_number,
                            "mime_type": mime_type,
                            "placement": "append_only",
                        },
                    )
                )
        return assets

    @staticmethod
    def _build_asset_publisher(
        afts_service: Optional["AftsService"],
    ) -> DocumentAssetPublisher:
        if afts_service is None:
            return NoOpDocumentAssetPublisher()
        return AftsDocumentAssetPublisher(afts_service)

    def _validate_published_assets(self, assets: list[DocumentAsset]) -> None:
        if self._asset_reference_mode == "local_path":
            return
        missing_remote_ids = [
            asset.asset_id
            for asset in assets
            if asset.local_path is not None and not asset.remote_id
        ]
        if missing_remote_ids:
            raise RuntimeError(
                "PPTX embedded image upload to AFTS failed: "
                + ", ".join(missing_remote_ids)
            )


def _load_presentation(file_path: Path) -> Any:
    if Presentation is None:
        raise RuntimeError("python_pptx provider requires python-pptx")
    return Presentation(str(file_path))


def _is_picture_shape(shape: Any) -> bool:
    if MSO_SHAPE_TYPE is not None:
        return getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE
    return hasattr(shape, "image")


def _get_shape_image_or_none(
    shape: Any,
    *,
    file_path: Path,
    slide_number: int,
) -> Any | None:
    """Return a real image payload while tolerating broken/non-image OOXML parts."""
    try:
        image = shape.image
        blob = image.blob
    except (AttributeError, KeyError, TypeError, ValueError) as exc:
        logger.warning(
            "ppt_document_service skip unsupported picture part | "
            "file_path=%s slide_number=%s shape_name=%s error=%s",
            file_path,
            slide_number,
            str(getattr(shape, "name", "") or ""),
            exc,
        )
        return None
    if not isinstance(blob, bytes):
        logger.warning(
            "ppt_document_service skip picture with invalid blob | "
            "file_path=%s slide_number=%s shape_name=%s blob_type=%s",
            file_path,
            slide_number,
            str(getattr(shape, "name", "") or ""),
            type(blob).__name__,
        )
        return None
    return image


def _image_extension(image: Any) -> str | None:
    blob = image.blob
    if blob.startswith(b"RIFF") and len(blob) >= 12 and blob[8:12] == b"WEBP":
        return ".webp"
    try:
        content_type = str(getattr(image, "content_type", "") or "").lower()
    except (AttributeError, KeyError, TypeError, ValueError):
        content_type = ""
    known_extension = _IMAGE_CONTENT_TYPES.get(content_type)
    if known_extension:
        return known_extension
    try:
        extension = str(getattr(image, "ext", "") or "").lstrip(".")
    except (AttributeError, KeyError, TypeError, ValueError):
        return None
    if not extension:
        return None
    return f".{extension}"


def _image_content_type(image: Any, extension: str) -> str:
    try:
        content_type = str(getattr(image, "content_type", "") or "").lower()
    except (AttributeError, KeyError, TypeError, ValueError):
        content_type = ""
    if content_type:
        return content_type
    normalized_extension = extension.lower()
    for known_content_type, known_extension in _IMAGE_CONTENT_TYPES.items():
        if known_extension == normalized_extension:
            return known_content_type
    return "application/octet-stream"


def _iter_shapes(shapes: Any):
    for shape in shapes:
        if _is_group_shape(shape):
            yield from _iter_shapes(shape.shapes)
            continue
        yield shape


def _is_group_shape(shape: Any) -> bool:
    if MSO_SHAPE_TYPE is not None:
        return getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP
    return hasattr(shape, "shapes")


def _escape_table_cell(value: Any) -> str:
    return str(value or "").strip().replace("|", "\\|").replace("\r", " ").replace("\n", "<br>")


def _parse_bool(value: Any, *, default: bool) -> bool:
    if value is None:
        return default
    if isinstance(value, bool):
        return value
    normalized = str(value).strip().lower()
    if normalized in {"1", "true", "yes", "on"}:
        return True
    if normalized in {"0", "false", "no", "off"}:
        return False
    raise ValueError(f"Invalid boolean value: {value}")
